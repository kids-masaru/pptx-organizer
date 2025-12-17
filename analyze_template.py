#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
テンプレートPPTXの目次スライドを分析するスクリプト（詳細版）
結果をファイルに出力
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

def analyze_slide_to_file(slide, slide_index, f):
    """スライドの構造を詳細に分析"""
    f.write(f"\n{'='*60}\n")
    f.write(f"スライド {slide_index + 1}\n")
    f.write('='*60 + '\n')
    
    for shape_idx, shape in enumerate(slide.shapes):
        f.write(f"\n  Shape {shape_idx}: {shape.shape_type.name}\n")
        f.write(f"    ID: {shape.shape_id}\n")
        f.write(f"    Name: {shape.name}\n")
        f.write(f"    Position: Left={shape.left.emu}, Top={shape.top.emu}, Width={shape.width.emu}, Height={shape.height.emu}\n")
        
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()[:500]
            text_display = text.replace('\n', '\\n')
            f.write(f"    HasTextFrame: True\n")
            f.write(f"    Text: '{text_display}'\n")
            f.write(f"    NumParagraphs: {len(shape.text_frame.paragraphs)}\n")
            
            for p_idx, para in enumerate(shape.text_frame.paragraphs[:5]):
                para_text = para.text[:150].replace('\n', '\\n')
                f.write(f"      Para{p_idx}: '{para_text}'\n")
                if para.font:
                    f.write(f"        Font: Bold={para.font.bold}, Size={para.font.size}\n")
                if para.level is not None:
                    f.write(f"        Level: {para.level}\n")

def main():
    pptx_path = "template.pptx"
    output_file = "template_analysis.txt"
    
    prs = Presentation(pptx_path)
    
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(f"テンプレートを分析: {pptx_path}\n")
        f.write(f"総スライド数: {len(prs.slides)}\n")
        
        # 最初の5スライドを詳細に分析
        for idx in range(min(5, len(prs.slides))):
            analyze_slide_to_file(prs.slides[idx], idx, f)
    
    print(f"分析結果を保存しました: {output_file}")

if __name__ == '__main__':
    main()
