#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPTX Organizer Script v4 (AI-Powered)
======================================
Gemini AIを使用してPDFの審査基準とPPTXスライドを
意味的にマッチングし、スライドを並べ替えるスクリプト。

Usage:
    python main.py <source_pdf> <master_pptx> [output_pptx]

Example:
    python main.py 審査基準表.pdf 【標準提案資料】2025-10-3.pptx output.pptx

Environment:
    GOOGLE_API_KEY: Gemini API Key (required)
"""

import sys
import os
import logging
import argparse
import re
import json
from pathlib import Path
from typing import List, Dict, Optional, Tuple

import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv

import google.generativeai as genai

# Load environment variables
load_dotenv()

# ============================================================================
# Logging Setup
# ============================================================================
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger(__name__)

# ============================================================================
# Gemini API Setup
# ============================================================================
def setup_gemini():
    """Gemini APIを初期化する"""
    api_key = os.getenv("GOOGLE_API_KEY")
    if not api_key:
        logger.error("GOOGLE_API_KEY が設定されていません。.env ファイルを確認してください。")
        sys.exit(1)
    
    genai.configure(api_key=api_key)
    
    # Gemini 2.5 Flash を使用
    model = genai.GenerativeModel("models/gemini-2.5-flash")
    logger.info("Gemini API 初期化完了 (gemini-2.5-flash)")
    return model


# ============================================================================
# Multi-Format File Extraction (PDF, Excel, Word, Image)
# ============================================================================
def detect_file_type(file_path: str) -> str:
    """ファイルタイプを検出"""
    ext = Path(file_path).suffix.lower()
    type_map = {
        '.pdf': 'pdf',
        '.xlsx': 'excel', '.xls': 'excel',
        '.docx': 'word', '.doc': 'word',
        '.png': 'image', '.jpg': 'image', '.jpeg': 'image',
    }
    return type_map.get(ext, 'unknown')


def extract_categories_from_pdf(pdf_path: str) -> List[Dict]:
    """
    審査基準表PDFから大項目・小項目を含む階層構造でカテゴリを抽出
    
    Returns:
        List[Dict]: 各要素は以下の構造
        {
            'No': int,
            'MainCategory': str,  # 大項目名
            'SubItems': List[str]  # 小項目のリスト
        }
    """
    logger.info(f"PDFを読み込み中: {pdf_path}")
    categories = []
    current_category = None
    
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                logger.info(f"  ページ {page_num} を処理中...")
                tables = page.extract_tables()
                
                for table in tables:
                    if not table or len(table) < 2:
                        continue
                    
                    for row in table:
                        if not row or len(row) < 2:
                            continue
                        
                        col0 = str(row[0]).strip() if row[0] else ""
                        col1 = str(row[1]).strip() if row[1] else ""
                        
                        # 大項目の検出（数字で始まる行）
                        if col0 and re.match(r'^\d+', col0):
                            no_match = re.match(r'^(\d+)', col0)
                            if no_match and col1:
                                # 前のカテゴリを保存
                                if current_category:
                                    categories.append(current_category)
                                
                                no = int(no_match.group(1))
                                # col1から大項目名を抽出（最初の行）
                                lines = col1.split('\n')
                                main_category = lines[0].strip()
                                
                                # 新しいカテゴリを開始
                                current_category = {
                                    'No': no,
                                    'MainCategory': main_category,
                                    'SubItems': []
                                }
                        
                        # 小項目の検出（col2に数字、col3に内容がある行）
                        if len(row) >= 4 and current_category:
                            col2 = str(row[2]).strip() if row[2] else ""
                            col3 = str(row[3]).strip() if row[3] else ""
                            
                            # col2 が小項目番号（1, 2, 3等）で col3 に内容がある場合
                            if col2 and re.match(r'^\d+$', col2) and col3:
                                # 既存の小項目と重複しないかチェック
                                sub_item = col3.split('\n')[0].strip()[:100]
                                if sub_item and sub_item not in current_category['SubItems']:
                                    current_category['SubItems'].append(sub_item)
                    
                # 最後のカテゴリを保存
                if current_category and current_category not in categories:
                    categories.append(current_category)
                    current_category = None
                            
    except Exception as e:
        logger.error(f"PDF読み込みエラー: {e}")
        raise
    
    # 重複を除去してソート
    seen_nos = set()
    unique_categories = []
    for cat in categories:
        if cat['No'] not in seen_nos:
            seen_nos.add(cat['No'])
            unique_categories.append(cat)
    
    unique_categories.sort(key=lambda x: x['No'])
    
    logger.info(f"抽出完了: {len(unique_categories)} 件のカテゴリ")
    for cat in unique_categories:
        logger.info(f"  No.{cat['No']:2d}: {cat['MainCategory']}")
        for sub in cat.get('SubItems', []):
            logger.info(f"       - {sub[:50]}...")
    
    return unique_categories



def extract_categories_from_excel(excel_path: str) -> List[Dict[str, str]]:
    """Excelファイルからカテゴリを抽出"""
    logger.info(f"Excelを読み込み中: {excel_path}")
    categories = []
    
    try:
        import openpyxl
        wb = openpyxl.load_workbook(excel_path)
        ws = wb.active
        
        for row in ws.iter_rows(values_only=True):
            if not row or len(row) < 2:
                continue
            
            col0 = str(row[0]).strip() if row[0] else ""
            col1 = str(row[1]).strip() if row[1] else ""
            
            if col0 and re.match(r'^\d+', col0):
                no_match = re.match(r'^(\d+)', col0)
                if no_match and col1:
                    no = int(no_match.group(1))
                    if not any(c['No'] == no for c in categories):
                        categories.append({'No': no, 'Category': col1.split('\n')[0].strip()})
    except Exception as e:
        logger.error(f"Excel読み込みエラー: {e}")
        raise
    
    categories.sort(key=lambda x: x['No'])
    logger.info(f"抽出完了: {len(categories)} 件のカテゴリ")
    return categories


def extract_categories_with_ai(model, file_path: str) -> List[Dict[str, str]]:
    """Gemini AIを使用してファイルからカテゴリを抽出（Word/Image対応）"""
    logger.info(f"AIでファイルを分析中: {file_path}")
    
    file_type = detect_file_type(file_path)
    
    if file_type == 'image':
        # 画像ファイルをアップロード
        uploaded_file = genai.upload_file(file_path)
        prompt_parts = [uploaded_file]
    else:
        # Word等はテキスト抽出
        try:
            import docx
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        except:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        prompt_parts = [text]
    
    prompt = """以下のファイルから審査基準のカテゴリ一覧を抽出してください。

出力形式（JSON）:
```json
[
  {"No": 1, "Category": "カテゴリ名"},
  {"No": 2, "Category": "カテゴリ名"}
]
```

番号順に並べてください。カテゴリ名は簡潔に（最初の1行程度）。
必ずJSON形式のみを出力してください。
"""
    
    try:
        response = model.generate_content([prompt] + prompt_parts)
        response_text = response.text.strip()
        
        if "```json" in response_text:
            response_text = response_text.split("```json")[1].split("```")[0].strip()
        elif "```" in response_text:
            response_text = response_text.split("```")[1].split("```")[0].strip()
        
        categories = json.loads(response_text)
        categories.sort(key=lambda x: x['No'])
        logger.info(f"AI抽出完了: {len(categories)} 件のカテゴリ")
        return categories
    except Exception as e:
        logger.error(f"AI抽出エラー: {e}")
        return []


def extract_categories(model, file_path: str) -> List[Dict[str, str]]:
    """ファイルタイプに応じてカテゴリを抽出（メイン関数）"""
    file_type = detect_file_type(file_path)
    logger.info(f"ファイルタイプ: {file_type}")
    
    if file_type == 'pdf':
        return extract_categories_from_pdf(file_path)
    elif file_type == 'excel':
        return extract_categories_from_excel(file_path)
    elif file_type in ('word', 'image'):
        return extract_categories_with_ai(model, file_path)
    else:
        logger.warning(f"未対応のファイル形式: {file_type}. AIで処理を試みます。")
        return extract_categories_with_ai(model, file_path)


# ============================================================================
# PPTX Utilities
# ============================================================================
def get_slide_title(slide) -> str:
    """スライドからタイトルを取得"""
    if slide.shapes.title:
        return slide.shapes.title.text.strip()
    return ""


def get_slide_first_text(slide) -> str:
    """スライドから最初のテキストを取得"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text
    return ""


def get_slide_full_content(slide) -> str:
    """
    スライドの全テキスト内容を取得（AIマッチング精度向上用）
    """
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and len(text) > 2:  # 2文字以上のテキストのみ
                texts.append(text)
    return "\n".join(texts)[:500]  # 最大500文字


def get_slide_groups(prs) -> List[Dict]:
    """スライドをグループ化（タイトル付きスライドを先頭に）"""
    groups = []
    current_group = None
    
    for idx, slide in enumerate(prs.slides):
        title = get_slide_title(slide)
        
        if title:
            if current_group:
                groups.append(current_group)
            current_group = {
                'title': title,
                'slides': [idx],
                'first_index': idx
            }
        else:
            if current_group:
                current_group['slides'].append(idx)
            else:
                first_text = get_slide_first_text(slide)
                current_group = {
                    'title': first_text[:50] if first_text else f"[Untitled {idx}]",
                    'slides': [idx],
                    'first_index': idx
                }
    
    if current_group:
        groups.append(current_group)
    
    return groups


# ============================================================================
# AI Matching with Gemini
# ============================================================================
def create_matching_with_ai(model, pdf_categories: List[Dict], pptx_groups: List[Dict]) -> Dict[int, int]:
    """
    Gemini AIを使用してPDFカテゴリとPPTXグループをマッチング。
    大項目・小項目とスライドの全テキスト内容を考慮してマッチング精度を向上。
    
    Returns:
        Dict[int, int]: {pdf_no: pptx_group_index} のマッピング
    """
    logger.info("")
    logger.info("=" * 60)
    logger.info("Gemini AI マッチング開始（精度向上版）")
    logger.info("=" * 60)
    
    # プロンプト用のデータを準備（階層構造を含める）
    pdf_entries = []
    for cat in pdf_categories:
        main_cat = cat.get('MainCategory', cat.get('Category', ''))
        sub_items = cat.get('SubItems', [])
        entry = f"PDF{cat['No']}: 【大項目】 {main_cat}"
        if sub_items:
            entry += f"\n  小項目: {', '.join(sub_items[:3])}"
        pdf_entries.append(entry)
    pdf_list = "\n".join(pdf_entries)
    
    # PPTXグループ情報（タイトル + 内容の要約）
    pptx_entries = []
    for i, g in enumerate(pptx_groups):
        content_summary = g.get('content', '')[:200] if g.get('content') else ''
        entry = f"PPTX{i}: {g['title']}"
        if content_summary:
            entry += f"\n  内容: {content_summary}..."
        pptx_entries.append(entry)
    pptx_list = "\n".join(pptx_entries)
    
    prompt = f"""あなたはドキュメント整理の専門家です。以下のタスクを実行してください。

## タスク
PDFの審査基準（大項目と小項目）と、PPTXのスライドグループを意味的にマッチングしてください。
**大項目だけでなく、小項目の内容も考慮して** 最も関連性の高いスライドグループを選んでください。

## PDFカテゴリ一覧（大項目と小項目）
{pdf_list}

## PPTXスライドグループ一覧（タイトルと内容）
{pptx_list}

## マッチングのルール
1. 大項目のテーマに最も近いスライドグループを選ぶ
2. 小項目の詳細内容も考慮して判断する
3. 表現が違っても同じトピックならマッチさせる
4. 1つのPPTXグループは1つのPDFカテゴリにのみマッチさせる

## 出力形式
JSON形式で出力。PDFのNo（数字）をキー、PPTXのインデックス（数字）を値とする。
マッチなしは-1。

例: {{"1": 3, "2": 5, "3": -1, "4": 7}}

必ずJSON形式のみを出力してください（説明は不要）。

出力:"""

    try:
        response = model.generate_content(prompt)
        response_text = response.text.strip()
        
        # JSONを抽出（マークダウンコードブロックに囲まれている場合の対応）
        if "```json" in response_text:
            response_text = response_text.split("```json")[1].split("```")[0].strip()
        elif "```" in response_text:
            response_text = response_text.split("```")[1].split("```")[0].strip()
        
        logger.info(f"AI応答: {response_text}")
        
        # JSONをパース
        mapping_raw = json.loads(response_text)
        
        # キーを整数に変換
        mapping = {}
        for pdf_no, pptx_idx in mapping_raw.items():
            pdf_no_int = int(pdf_no)
            if pptx_idx >= 0:
                mapping[pdf_no_int] = pptx_idx
        
        logger.info(f"マッチング結果: {len(mapping)} 件")
        
        return mapping
        
    except json.JSONDecodeError as e:
        logger.error(f"JSON解析エラー: {e}")
        logger.error(f"応答テキスト: {response_text}")
        return {}
    except Exception as e:
        logger.error(f"AI マッチングエラー: {e}")
        return {}


# ============================================================================
# Main Processing
# ============================================================================
def populate_toc(prs, categories: List[Dict], toc_slide_index: int = 1):
    """
    目次スライドに審査基準カテゴリを階層構造で入力
    大項目は太字、小項目はインデントして表示
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    
    logger.info(f"目次スライド（インデックス {toc_slide_index}）にカテゴリを入力中...")
    
    try:
        toc_slide = prs.slides[toc_slide_index]
        
        # テンプレートの目次用テキストボックスを探す（最も大きいテキストフレーム）
        target_shape = None
        max_area = 0
        
        for shape in toc_slide.shapes:
            if shape.has_text_frame:
                area = shape.width * shape.height
                existing_text = shape.text_frame.text.strip()
                # タイトルやページ番号を除外（小さいテキストや数字のみ）
                if len(existing_text) > 10 and area > max_area:
                    max_area = area
                    target_shape = shape
        
        if not target_shape:
            logger.warning("  目次用のテキストフレームが見つかりませんでした")
            return False
        
        # テキストフレームをクリア
        tf = target_shape.text_frame
        tf.clear()
        
        # 階層構造の目次を作成
        for idx, cat in enumerate(categories):
            # 大項目
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # 大項目の番号とタイトル
            main_text = cat.get('MainCategory', cat.get('Category', ''))
            p.text = f"{cat['No']}. {main_text}"
            
            # 大項目のフォーマット（太字）
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(11)
            if not p.runs:
                # テキストが直接設定された場合
                p.font.bold = True
                p.font.size = Pt(11)
            
            # 小項目を追加
            sub_items = cat.get('SubItems', [])
            for sub_item in sub_items[:3]:  # 最大3件の小項目を表示
                sub_p = tf.add_paragraph()
                sub_p.text = f"  ・ {sub_item[:40]}"  # インデント + 中点
                sub_p.level = 1
                
                # 小項目のフォーマット（通常）
                for run in sub_p.runs:
                    run.font.bold = False
                    run.font.size = Pt(9)
                if not sub_p.runs:
                    sub_p.font.bold = False
                    sub_p.font.size = Pt(9)
        
        logger.info(f"  目次を更新しました: {len(categories)} 項目（階層構造）")
        return True
        
    except Exception as e:
        logger.error(f"  目次入力エラー: {e}")
        import traceback
        traceback.print_exc()
        return False


def update_slide_title(slide, new_title: str):
    """スライドのタイトルを更新"""
    try:
        if slide.shapes.title:
            slide.shapes.title.text = new_title
            return True
    except:
        pass
    return False


def process_pptx(model, pdf_categories: List[Dict], pptx_path: str, output_path: str):
    """PDFカテゴリに基づいてPPTXを処理（表紙・目次を固定）"""
    logger.info(f"PPTXを読み込み中: {pptx_path}")
    prs = Presentation(pptx_path)
    
    total_slides = len(prs.slides)
    logger.info(f"スライド総数: {total_slides}")
    
    # 表紙（0）と目次（1）を固定するため、スライド2以降をグループ化
    FIXED_SLIDES = 2  # 表紙と目次
    
    if total_slides <= FIXED_SLIDES:
        logger.error("スライドが少なすぎます。")
        return
    
    # 目次に審査基準カテゴリを入力
    logger.info("")
    logger.info("=" * 60)
    logger.info("目次スライドの更新")
    logger.info("=" * 60)
    populate_toc(prs, pdf_categories, toc_slide_index=1)
    
    # スライド2以降をグループ化
    groups = []
    current_group = None
    
    for idx in range(FIXED_SLIDES, total_slides):
        slide = prs.slides[idx]
        title = get_slide_title(slide)
        content = get_slide_full_content(slide)  # スライドの全内容を取得
        
        if title:
            if current_group:
                groups.append(current_group)
            current_group = {
                'title': title,
                'slides': [idx],
                'first_index': idx,
                'content': content  # AIマッチング用にコンテンツを追加
            }
        else:
            if current_group:
                current_group['slides'].append(idx)
                # コンテンツを累積
                current_group['content'] = (current_group.get('content', '') + '\n' + content)[:500]
            else:
                first_text = get_slide_first_text(slide)
                current_group = {
                    'title': first_text[:50] if first_text else f"[Untitled {idx}]",
                    'slides': [idx],
                    'first_index': idx,
                    'content': content
                }
    
    if current_group:
        groups.append(current_group)
    
    logger.info(f"コンテンツスライドグループ数: {len(groups)}")
    
    for i, g in enumerate(groups):
        logger.info(f"  Group {i}: '{g['title'][:50]}...' - Slides {[idx+1 for idx in g['slides']]}")
    
    # AIでマッチング
    mapping = create_matching_with_ai(model, pdf_categories, groups)
    
    if not mapping:
        logger.error("マッチングに失敗しました。")
        return
    
    # マッチング結果を表示
    logger.info("")
    logger.info("=" * 60)
    logger.info("マッチング結果詳細")
    logger.info("=" * 60)
    
    used_groups = set()
    matched_list = []
    
    for cat in pdf_categories:
        pdf_no = cat['No']
        main_cat = cat.get('MainCategory', cat.get('Category', ''))
        if pdf_no in mapping:
            pptx_idx = mapping[pdf_no]
            if pptx_idx < len(groups):
                group = groups[pptx_idx]
                logger.info(f"  ✓ PDF[{pdf_no}] '{main_cat[:30]}...'")
                logger.info(f"    → PPTX '{group['title'][:40]}...' ({len(group['slides'])} slides)")
                matched_list.append((pdf_no, main_cat, group))
                used_groups.add(pptx_idx)
        else:
            logger.info(f"  ✗ PDF[{pdf_no}] '{main_cat[:30]}...' - マッチなし")
    
    # 未使用グループ
    unused_groups = [g for i, g in enumerate(groups) if i not in used_groups]
    
    logger.info("")
    logger.info("=" * 60)
    logger.info("スライド再構成開始")
    logger.info("=" * 60)
    
    # 新しい順序を構築（表紙・目次は固定）
    new_order = list(range(FIXED_SLIDES))  # [0, 1] = 表紙と目次
    
    # マッチしたグループをPDF順に配置
    matched_list.sort(key=lambda x: x[0])
    for pdf_no, category_name, group in matched_list:
        # 大項目スライドのタイトルを更新
        first_slide_idx = group['slides'][0]
        new_title = f"{pdf_no}. {category_name}"
        if update_slide_title(prs.slides[first_slide_idx], new_title):
            logger.info(f"  タイトル更新: '{new_title}'")
        
        for slide_idx in group['slides']:
            new_order.append(slide_idx)
        logger.info(f"  配置 No.{pdf_no}: '{category_name[:40]}...' ({len(group['slides'])} slides)")
    
    # 未使用グループを末尾に配置
    if unused_groups:
        logger.info(f"  --- 以下、未使用スライド ---")
        for g in unused_groups:
            for slide_idx in g['slides']:
                new_order.append(slide_idx)
            logger.info(f"  末尾: '{g['title'][:40]}...' ({len(g['slides'])} slides)")
    
    # XMLレベルでスライドを並べ替え
    xml_slides = prs.slides._sldIdLst
    original_slides = list(xml_slides)
    
    while len(xml_slides) > 0:
        xml_slides.remove(xml_slides[0])
    
    for idx in new_order:
        xml_slides.append(original_slides[idx])
    
    # 保存
    logger.info("")
    logger.info(f"保存中: {output_path}")
    prs.save(output_path)
    logger.info("完了!")
    
    # サマリー
    logger.info("")
    logger.info("=" * 60)
    logger.info("処理サマリー")
    logger.info("=" * 60)
    logger.info(f"  マッチしたグループ: {len(matched_list)}")
    logger.info(f"  未使用グループ: {len(unused_groups)}")
    logger.info(f"  出力ファイル: {output_path}")


# ============================================================================
# Main Entry Point
# ============================================================================
def main():
    parser = argparse.ArgumentParser(
        description='Gemini AIを使用してPPTXスライドを並べ替え（表紙・目次を固定、タイトル自動更新）'
    )
    parser.add_argument('source_file', 
                        help='審査基準ファイル（PDF, Excel, Word, 画像に対応）')
    parser.add_argument('master_pptx', help='編集対象のPPTXファイル')
    parser.add_argument('output_pptx', nargs='?', default=None,
                        help='出力PPTXファイル（省略時は {master}_output.pptx）')
    
    args = parser.parse_args()
    
    source_path = Path(args.source_file)
    pptx_path = Path(args.master_pptx)
    
    if not source_path.exists():
        logger.error(f"審査基準ファイルが見つかりません: {source_path}")
        sys.exit(1)
    
    if not pptx_path.exists():
        logger.error(f"PPTXファイルが見つかりません: {pptx_path}")
        sys.exit(1)
    
    if args.output_pptx:
        output_path = Path(args.output_pptx)
    else:
        output_path = pptx_path.parent / f"{pptx_path.stem}_output{pptx_path.suffix}"
    
    logger.info("=" * 60)
    logger.info("PPTX Organizer v5 (AI-Powered) - スライド自動整理ツール")
    logger.info("=" * 60)
    logger.info(f"入力ファイル: {source_path} ({detect_file_type(str(source_path))})")
    logger.info(f"入力PPTX: {pptx_path}")
    logger.info(f"出力PPTX: {output_path}")
    logger.info("")
    
    try:
        # Gemini API初期化
        model = setup_gemini()
        
        # ファイル形式に応じてカテゴリ抽出
        categories = extract_categories(model, str(source_path))
        
        if not categories:
            logger.error("審査基準からカテゴリを抽出できませんでした。")
            sys.exit(1)
        
        # PPTX 処理
        process_pptx(model, categories, str(pptx_path), str(output_path))
        
    except Exception as e:
        logger.error(f"処理中にエラーが発生しました: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
