"""
PPTX Organizer - Streamlit Web Application
==========================================
審査基準に基づいてPowerPointスライドを自動整理するWebアプリ

Features:
- 審査基準ファイルアップロード (PDF/Excel/Word/画像)
- PPTXテンプレートアップロード
- AI自動マッチング＆並べ替え
- 結果PPTXダウンロード
"""

import streamlit as st
import tempfile
import os
import io
import re
import json
from pathlib import Path

import pdfplumber
from pptx import Presentation
import google.generativeai as genai

# ============================================================================
# Page Config
# ============================================================================
st.set_page_config(
    page_title="PPTX Organizer",
    page_icon="📊",
    layout="wide"
)

# ============================================================================
# Custom CSS
# ============================================================================
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        font-weight: bold;
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 1rem;
    }
    .upload-section {
        background: #f8f9fa;
        padding: 1.5rem;
        border-radius: 10px;
        margin-bottom: 1rem;
    }
    .status-box {
        padding: 1rem;
        border-radius: 8px;
        margin: 0.5rem 0;
    }
    .success-box {
        background-color: #d4edda;
        border-left: 4px solid #28a745;
    }
</style>
""", unsafe_allow_html=True)

# ============================================================================
# Gemini API Setup
# ============================================================================
def setup_gemini():
    """Gemini APIを初期化"""
    api_key = None
    
    # Streamlit Secretsから取得
    if "GOOGLE_API_KEY" in st.secrets:
        api_key = st.secrets["GOOGLE_API_KEY"]
    else:
        api_key = os.getenv("GOOGLE_API_KEY")
    
    if not api_key:
        st.error("⚠️ GOOGLE_API_KEY が設定されていません")
        st.stop()
    
    genai.configure(api_key=api_key)
    return genai.GenerativeModel("models/gemini-2.5-flash")

# ============================================================================
# File Type Detection
# ============================================================================
def detect_file_type(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    type_map = {
        '.pdf': 'pdf',
        '.xlsx': 'excel', '.xls': 'excel',
        '.docx': 'word', '.doc': 'word',
        '.png': 'image', '.jpg': 'image', '.jpeg': 'image',
    }
    return type_map.get(ext, 'unknown')

# ============================================================================
# Category Extraction Functions
# ============================================================================
def extract_categories_from_pdf(file_bytes) -> list:
    """PDFからカテゴリを抽出"""
    categories = []
    
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                if not table or len(table) < 2:
                    continue
                for row in table:
                    if not row or len(row) < 2:
                        continue
                    col0 = str(row[0]).strip() if row[0] else ""
                    col1 = str(row[1]).strip() if row[1] else ""
                    
                    if col0 and re.match(r'^\d+', col0):
                        no_match = re.match(r'^(\d+)', col0)
                        if no_match and col1:
                            no = int(no_match.group(1))
                            category = col1.split('\n')[0].strip()
                            if not any(c['No'] == no for c in categories):
                                categories.append({'No': no, 'Category': category})
    
    categories.sort(key=lambda x: x['No'])
    return categories


def extract_categories_from_excel(file_bytes) -> list:
    """Excelからカテゴリを抽出"""
    import openpyxl
    categories = []
    
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
    ws = wb.active
    
    for row in ws.iter_rows(values_only=True):
        if not row or len(row) < 2:
            continue
        col0 = str(row[0]).strip() if row[0] else ""
        col1 = str(row[1]).strip() if row[1] else ""
        
        if col0 and re.match(r'^\d+', col0):
            no_match = re.match(r'^(\d+)', col0)
            if no_match and col1:
                no = int(no_match.group(1))
                if not any(c['No'] == no for c in categories):
                    categories.append({'No': no, 'Category': col1.split('\n')[0].strip()})
    
    categories.sort(key=lambda x: x['No'])
    return categories


def extract_categories_with_ai(model, file_bytes, file_type: str) -> list:
    """AIでカテゴリを抽出（Word/画像）"""
    
    if file_type == 'image':
        # 一時ファイルに保存してアップロード
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        uploaded_file = genai.upload_file(tmp_path)
        prompt_parts = [uploaded_file]
        os.unlink(tmp_path)
    else:
        # Wordはテキスト抽出
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            text = "\n".join([para.text for para in doc.paragraphs])
        except:
            text = file_bytes.decode('utf-8', errors='ignore')
        prompt_parts = [text]
    
    prompt = """以下のファイルから審査基準のカテゴリ一覧を抽出してください。

出力形式（JSON）:
```json
[
  {"No": 1, "Category": "カテゴリ名"},
  {"No": 2, "Category": "カテゴリ名"}
]
```

番号順に並べてください。必ずJSON形式のみを出力してください。
"""
    
    response = model.generate_content([prompt] + prompt_parts)
    response_text = response.text.strip()
    
    if "```json" in response_text:
        response_text = response_text.split("```json")[1].split("```")[0].strip()
    elif "```" in response_text:
        response_text = response_text.split("```")[1].split("```")[0].strip()
    
    categories = json.loads(response_text)
    categories.sort(key=lambda x: x['No'])
    return categories


def extract_categories(model, file_bytes, filename: str) -> list:
    """ファイル形式に応じてカテゴリを抽出"""
    file_type = detect_file_type(filename)
    
    if file_type == 'pdf':
        return extract_categories_from_pdf(file_bytes)
    elif file_type == 'excel':
        return extract_categories_from_excel(file_bytes)
    elif file_type in ('word', 'image'):
        return extract_categories_with_ai(model, file_bytes, file_type)
    else:
        return extract_categories_with_ai(model, file_bytes, 'word')

# ============================================================================
# PPTX Processing Functions
# ============================================================================
def get_slide_title(slide) -> str:
    if slide.shapes.title:
        return slide.shapes.title.text.strip()
    return ""


def get_slide_first_text(slide) -> str:
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text
    return ""


def populate_toc(prs, categories, toc_slide_index=1):
    """目次スライドにカテゴリを入力"""
    try:
        toc_slide = prs.slides[toc_slide_index]
        toc_text = "\n".join([f"{cat['No']}. {cat['Category']}" for cat in categories])
        
        for shape in toc_slide.shapes:
            if shape.has_text_frame:
                existing_text = shape.text_frame.text.strip()
                if existing_text and (existing_text.isdigit() or re.match(r'^[\d\s\n]+$', existing_text)):
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = toc_text
                    return True
        return False
    except:
        return False


def update_slide_title(slide, new_title: str):
    try:
        if slide.shapes.title:
            slide.shapes.title.text = new_title
            return True
    except:
        pass
    return False


def create_matching_with_ai(model, categories, groups) -> dict:
    """AIでマッチング"""
    cat_list = "\n".join([f"CAT{cat['No']}: {cat['Category']}" for cat in categories])
    grp_list = "\n".join([f"GRP{i}: {g['title']}" for i, g in enumerate(groups)])
    
    prompt = f"""審査基準カテゴリとPPTXスライドグループをマッチングしてください。

## カテゴリ一覧
{cat_list}

## スライドグループ一覧
{grp_list}

## 出力形式
JSON形式で。カテゴリNoをキー、グループインデックスを値。マッチなしは-1。
例: {{"1": 3, "2": 5, "3": -1}}

必ずJSON形式のみを出力してください。
"""
    
    response = model.generate_content(prompt)
    response_text = response.text.strip()
    
    if "```json" in response_text:
        response_text = response_text.split("```json")[1].split("```")[0].strip()
    elif "```" in response_text:
        response_text = response_text.split("```")[1].split("```")[0].strip()
    
    mapping_raw = json.loads(response_text)
    mapping = {}
    for k, v in mapping_raw.items():
        if int(v) >= 0:
            mapping[int(k)] = int(v)
    return mapping


def process_pptx(model, categories, pptx_bytes, progress_callback=None) -> bytes:
    """PPTXを処理して並べ替え"""
    prs = Presentation(io.BytesIO(pptx_bytes))
    total_slides = len(prs.slides)
    
    FIXED_SLIDES = 2  # 表紙と目次
    
    if total_slides <= FIXED_SLIDES:
        raise ValueError("スライドが少なすぎます")
    
    # 目次更新
    if progress_callback:
        progress_callback(0.1, "目次を更新中...")
    populate_toc(prs, categories, toc_slide_index=1)
    
    # スライド2以降をグループ化
    if progress_callback:
        progress_callback(0.2, "スライドをグループ化中...")
    
    groups = []
    current_group = None
    
    for idx in range(FIXED_SLIDES, total_slides):
        slide = prs.slides[idx]
        title = get_slide_title(slide)
        
        if title:
            if current_group:
                groups.append(current_group)
            current_group = {'title': title, 'slides': [idx], 'first_index': idx}
        else:
            if current_group:
                current_group['slides'].append(idx)
            else:
                first_text = get_slide_first_text(slide)
                current_group = {
                    'title': first_text[:50] if first_text else f"[Untitled {idx}]",
                    'slides': [idx], 'first_index': idx
                }
    
    if current_group:
        groups.append(current_group)
    
    # AIマッチング
    if progress_callback:
        progress_callback(0.4, "AIでマッチング中...")
    
    mapping = create_matching_with_ai(model, categories, groups)
    
    # マッチング結果を整理
    if progress_callback:
        progress_callback(0.6, "スライドを並べ替え中...")
    
    used_groups = set()
    matched_list = []
    
    for cat in categories:
        pdf_no = cat['No']
        if pdf_no in mapping:
            pptx_idx = mapping[pdf_no]
            if pptx_idx < len(groups):
                matched_list.append((pdf_no, cat['Category'], groups[pptx_idx]))
                used_groups.add(pptx_idx)
    
    unused_groups = [g for i, g in enumerate(groups) if i not in used_groups]
    
    # 新しい順序を構築
    new_order = list(range(FIXED_SLIDES))
    
    matched_list.sort(key=lambda x: x[0])
    for pdf_no, category_name, group in matched_list:
        # タイトル更新
        first_slide_idx = group['slides'][0]
        new_title = f"{pdf_no}. {category_name}"
        update_slide_title(prs.slides[first_slide_idx], new_title)
        
        for slide_idx in group['slides']:
            new_order.append(slide_idx)
    
    for g in unused_groups:
        for slide_idx in g['slides']:
            new_order.append(slide_idx)
    
    # XMLレベルで並べ替え
    if progress_callback:
        progress_callback(0.8, "ファイルを生成中...")
    
    xml_slides = prs.slides._sldIdLst
    original_slides = list(xml_slides)
    
    while len(xml_slides) > 0:
        xml_slides.remove(xml_slides[0])
    
    for idx in new_order:
        xml_slides.append(original_slides[idx])
    
    # バイトに変換
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    if progress_callback:
        progress_callback(1.0, "完了！")
    
    return output.read(), len(matched_list), len(unused_groups)

# ============================================================================
# Template Management
# ============================================================================
TEMPLATE_PATH = Path(__file__).parent / "template.pptx"

def get_saved_template() -> bytes:
    """保存されたテンプレートを取得"""
    if TEMPLATE_PATH.exists():
        return TEMPLATE_PATH.read_bytes()
    return None

def save_template(file_bytes: bytes):
    """テンプレートを保存"""
    try:
        TEMPLATE_PATH.write_bytes(file_bytes)
        return True
    except:
        return False

# ============================================================================
# Main UI
# ============================================================================
st.markdown('<h1 class="main-header">📊 PPTX Organizer</h1>', unsafe_allow_html=True)
st.caption("審査基準に基づいてPowerPointスライドを自動整理")

# サイドバー
with st.sidebar:
    st.header("📋 使い方")
    st.markdown("""
    1. **審査基準ファイル**をアップロード
       - PDF / Excel / Word / 画像
    2. **処理開始**ボタンをクリック
    3. 完成したPPTXを**ダウンロード**
    """)
    
    st.markdown("---")
    
    # テンプレート設定（折りたたみ）
    with st.expander("⚙️ テンプレート更新", expanded=False):
        saved_template = get_saved_template()
        if saved_template:
            st.success(f"✅ 設定済み（{len(saved_template) / 1024 / 1024:.1f} MB）")
        else:
            st.warning("⚠️ テンプレートがありません")
        
        st.caption("新しいテンプレートに変更する場合のみ使用")
        template_upload = st.file_uploader(
            "新しいテンプレート",
            type=['pptx'],
            key="template_upload",
            label_visibility="collapsed"
        )
        
        if template_upload:
            if st.button("💾 更新を保存", use_container_width=True):
                template_bytes = template_upload.read()
                if save_template(template_bytes):
                    st.success("✅ 更新しました！")
                    st.rerun()
                else:
                    st.error("保存に失敗しました")

# メインエリア
st.subheader("📁 審査基準ファイル")
criteria_file = st.file_uploader(
    "審査基準をアップロード（PDF / Excel / Word / 画像）",
    type=['pdf', 'xlsx', 'xls', 'docx', 'doc', 'png', 'jpg', 'jpeg'],
    key="criteria"
)
if criteria_file:
    file_type = detect_file_type(criteria_file.name)
    st.success(f"✅ {criteria_file.name} ({file_type})")

# テンプレート状態表示
st.markdown("---")
template_to_use = get_saved_template()

if template_to_use:
    st.info("📊 保存済みテンプレートを使用します（サイドバーで変更可能）")
else:
    st.warning("⚠️ テンプレートがありません。サイドバーからアップロードしてください。")

# 処理ボタン
st.markdown("---")

if criteria_file and template_to_use:
    if st.button("🚀 処理開始", type="primary", use_container_width=True):
        try:
            model = setup_gemini()
            
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            def update_progress(value, text):
                progress_bar.progress(value)
                status_text.text(text)
            
            # カテゴリ抽出
            update_progress(0.05, "審査基準を分析中...")
            criteria_bytes = criteria_file.read()
            categories = extract_categories(model, criteria_bytes, criteria_file.name)
            
            if not categories:
                st.error("審査基準からカテゴリを抽出できませんでした")
                st.stop()
            
            st.info(f"📋 {len(categories)} 件のカテゴリを抽出しました")
            
            # PPTX処理
            result_bytes, matched_count, unused_count = process_pptx(
                model, categories, template_to_use, update_progress
            )
            
            # 結果表示
            st.success(f"✅ 処理完了！ マッチ: {matched_count}件 / 未使用: {unused_count}件")
            
            # ダウンロードボタン
            output_filename = f"organized_{criteria_file.name.split('.')[0]}.pptx"
            st.download_button(
                label="📥 完成PPTXをダウンロード",
                data=result_bytes,
                file_name=output_filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary",
                use_container_width=True
            )
            
            # カテゴリ一覧表示
            with st.expander("📋 抽出されたカテゴリ一覧"):
                for cat in categories:
                    st.write(f"{cat['No']}. {cat['Category']}")
                    
        except Exception as e:
            st.error(f"エラーが発生しました: {e}")
            import traceback
            st.code(traceback.format_exc())
elif not template_to_use:
    st.info("👈 サイドバーからテンプレートをアップロードしてください")
else:
    st.info("👆 審査基準ファイルをアップロードしてください")

# フッター
st.markdown("---")
st.caption("PPTX Organizer v5 | Powered by Google Gemini")

