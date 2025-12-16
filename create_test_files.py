#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
テスト用のサンプルPDFとPPTXを生成するスクリプト
"""

from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

from pptx import Presentation
from pptx.util import Inches, Pt

import os

# 出力ディレクトリ
OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

def create_sample_pdf():
    """テスト用PDFを作成"""
    pdf_path = os.path.join(OUTPUT_DIR, "test_instructions.pdf")
    
    doc = SimpleDocTemplate(pdf_path, pagesize=A4)
    elements = []
    
    styles = getSampleStyleSheet()
    
    # タイトル
    elements.append(Paragraph("Slide Reorganization Instructions", styles['Title']))
    elements.append(Paragraph(" ", styles['Normal']))
    
    # テーブルデータ
    data = [
        ['No', 'Old Title', 'New Title'],
        ['1', 'Introduction', '1. Introduction (Updated)'],
        ['2', 'Overview', '2. Project Overview'],
        ['3', 'Features', '3. Key Features'],
        ['4', 'New Section', '4. New Section'],  # これはPPTに存在しない
        ['5', 'Summary', '5. Conclusion'],
    ]
    
    table = Table(data, colWidths=[50, 200, 200])
    table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ]))
    
    elements.append(table)
    
    doc.build(elements)
    print(f"Created: {pdf_path}")
    return pdf_path


def create_sample_pptx():
    """テスト用PPTXを作成"""
    pptx_path = os.path.join(OUTPUT_DIR, "test_master.pptx")
    
    prs = Presentation()
    
    # スライドを追加
    slides_data = [
        "Introduction",           # マッチする予定
        "Overview",               # マッチする予定
        "Features",               # マッチする予定  
        "Summary",                # マッチする予定
        "Appendix",               # 未使用になる予定
        "Old Data",               # 未使用になる予定
    ]
    
    for title in slides_data:
        slide_layout = prs.slide_layouts[5]  # blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトルを追加
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(32)
    
    prs.save(pptx_path)
    print(f"Created: {pptx_path}")
    return pptx_path


if __name__ == '__main__':
    print("Generating test files...")
    create_sample_pdf()
    create_sample_pptx()
    print("Done!")
