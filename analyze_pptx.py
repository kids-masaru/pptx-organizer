#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""PPTXスライド構造を解析するスクリプト"""

from pptx import Presentation

prs = Presentation('【標準提案資料】2025-10-3.pptx')

print(f"Total slides: {len(prs.slides)}")
print("\n=== スライド一覧 ===")
for i, slide in enumerate(prs.slides):
    if slide.shapes.title:
        title = slide.shapes.title.text.strip()
    else:
        # タイトルが無い場合、最初のテキストを探す
        title = "[No Title]"
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    title = f"[No Title] First text: {text[:50]}..."
                    break
    print(f"Slide {i+1:3d}: {title}")
